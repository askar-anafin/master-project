import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Initialize presentation
    prs = Presentation()
    
    # Set aspect ratio to 16:9 (standard widescreen)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Design System colors (Matching the premium aesthetics of the example)
    CREAM = RGBColor(245, 244, 240)        # Warm cream background
    SLATE = RGBColor(15, 23, 42)          # Primary text (Navy-slate)
    RUST = RGBColor(184, 55, 38)          # Primary Accent (Rust Red)
    TEAL = RGBColor(38, 90, 84)           # Supporting Accent 1 (Teal)
    PEACH = RGBColor(230, 204, 184)        # Supporting Accent 2 (Peach)
    CARD_BG = RGBColor(248, 246, 240)      # Light cream fill for cards
    MUTED_GRAY = RGBColor(148, 163, 184)   # Divider lines & muted text
    
    FONT_SERIF = "Georgia"
    FONT_SANS = "Segoe UI"
    
    # Backward compatibility aliases
    WHITE = CREAM
    BLACK = SLATE
    DARK_SLATE = SLATE
    GRAY_LINE = MUTED_GRAY
    LIGHT_GRAY = CARD_BG
    FONT_TITLE = FONT_SERIF
    FONT_BODY = FONT_SANS

    # Helper function to apply slide background (warm cream)
    def apply_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = CREAM

    # Map titles to section categories for top-left header label
    CATEGORIES_MAP = {
        "Contents of Thesis": "CONTENTS",
        "Relevance of Thesis Topic": "RELEVANCE",
        "Research Aim": "RESEARCH AIM",
        "Objectives": "OBJECTIVES",
        "Literature Review": "LITERATURE REVIEW",
        "Novelty of the Research": "NOVELTY OF THE RESEARCH",
        "Problem Statement: Input/Output Tensors": "PROBLEM STATEMENT: INPUT/OUTPUT",
        "ECG Signal Data Table (12 Leads)": "PROBLEM STATEMENT: INPUT/OUTPUT",
        "Initial Dataset Overview": "COLLECTION, ANALYSIS AND PROCESSING OF INPUT DATA",
        "Zero-Phase Signal Preprocessing": "COLLECTION, ANALYSIS AND PROCESSING OF INPUT DATA",
        "Diagnostic Label Filtering": "COLLECTION, ANALYSIS AND PROCESSING OF INPUT DATA",
        "Baseline Models: CNN & ViT": "MODELS AND METHODS FOR SOLVING THE PROBLEM",
        "Spatio-Temporal ReGE GNN": "MODELS AND METHODS FOR SOLVING THE PROBLEM",
        "Spatio-Temporal Hybrid Model": "MODELS AND METHODS FOR SOLVING THE PROBLEM",
        "Core Outcomes: ST-ReGE & Hybrid": "MODELS AND METHODS FOR SOLVING THE PROBLEM",
        "Ensemble & Boundary Calibration": "MODELS AND METHODS FOR SOLVING THE PROBLEM",
        "Why Multi-Label Evaluation Metrics Matter": "MODELS AND METHODS FOR SOLVING THE PROBLEM",
        "Quantitative Model Benchmarks": "RESULTS",
        "Model Interpretability (1D Grad-CAM)": "RESULTS",
        "Clinical Error Analysis (Confusion Matrix)": "RESULTS",
        "FastAPI Diagnostic Triage Dashboard": "RESULTS",
        "Conclusion": "CONCLUSION",
        "Future Directions": "CONCLUSION",
        "Scholarly Publication": "RESULTS"
    }

    # Helper function to add a clean header and footer
    def add_slide_header(slide, title_text, slide_number, total_slides=22):
        # Determine category text from map
        category_text = CATEGORIES_MAP.get(title_text, "THESIS DEFENSE")
        
        # 1. Category text (top left)
        txBox_cat = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8.0), Inches(0.4))
        tf_cat = txBox_cat.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = FONT_SANS
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = RUST
        p_cat.alignment = PP_ALIGN.LEFT
        
        # 2. Page number (top right)
        slide_num_str = f"{slide_number:02d}"
        txBox_num = slide.shapes.add_textbox(Inches(11.833), Inches(0.1), Inches(1.0), Inches(0.5))
        tf_num = txBox_num.text_frame
        tf_num.word_wrap = True
        tf_num.margin_left = tf_num.margin_right = tf_num.margin_top = tf_num.margin_bottom = 0
        p_num = tf_num.paragraphs[0]
        p_num.text = slide_num_str
        p_num.font.name = FONT_SERIF
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = RUST
        p_num.alignment = PP_ALIGN.RIGHT
        
        # 3. Thin line separator
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.6), Inches(12.333), Inches(0.015)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = SLATE
        line.line.fill.background() # No border
        
        # 4. Slide title (below line)
        txBox_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.72), Inches(12.333), Inches(0.6))
        tf_title = txBox_title.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = FONT_SERIF
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = SLATE
        p_title.alignment = PP_ALIGN.LEFT

        # Footer text (very minimal and elegant)
        txBox_footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.15), Inches(12.333), Inches(0.25))
        tf_f = txBox_footer.text_frame
        tf_f.margin_left = tf_f.margin_right = tf_f.margin_top = tf_f.margin_bottom = 0
        p_f = tf_f.paragraphs[0]
        p_f.text = f"Askar Anafin | Master Thesis Defense | June 2026"
        p_f.font.name = FONT_SANS
        p_f.font.size = Pt(10)
        p_f.font.color.rgb = MUTED_GRAY
        p_f.alignment = PP_ALIGN.LEFT

    # Helper function to add a standard text box with formatting and custom bullets
    def add_textbox(slide, left, top, width, height, text_list, size=17, bold_first_n_words=0, bullet=True):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        for i, item in enumerate(text_list):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.space_after = Pt(8)
            p.font.name = FONT_SANS
            p.font.size = Pt(size)
            
            # Bullet prefix (Triangle ▸ in Rust color)
            if bullet and item.strip():
                run_bullet = p.add_run()
                run_bullet.text = "▸ "
                run_bullet.font.bold = True
                run_bullet.font.color.rgb = RUST
                
            # Parse bold text (colon-delimited)
            if ":" in item:
                parts = item.split(":", 1)
                run1 = p.add_run()
                run1.text = parts[0] + ":"
                run1.font.bold = True
                run1.font.color.rgb = SLATE
                
                run2 = p.add_run()
                run2.text = parts[1]
                run2.font.bold = False
                run2.font.color.rgb = SLATE
            else:
                run = p.add_run()
                run.text = item
                run.font.bold = False
                run.font.color.rgb = SLATE
                
        return txBox

    # Helper to add a clean card/block with thin border and top accent line
    def add_card(slide, left, top, width, height, title, body_paragraphs, color=None, border_color=None, font_size=15, accent_color=RUST):
        # Resolve colors to our design system
        card_fill = color if isinstance(color, RGBColor) else CARD_BG
        
        # 1. Main Card Shape (Light Alabaster Cream rounded rectangle)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = card_fill
        card.line.color.rgb = MUTED_GRAY
        card.line.width = Pt(0.75)
        
        # 2. Top Accent Line (Solid accent color bar)
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.08))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_color
        accent_bar.line.fill.background() # No border
        
        # 3. Inside Text Frame
        txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.3))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        if title:
            p_title = tf.paragraphs[0]
            p_title.text = title
            p_title.font.name = FONT_SERIF
            p_title.font.size = Pt(16)
            p_title.font.bold = True
            p_title.font.color.rgb = SLATE
            p_title.space_after = Pt(8)
            
        for i, para in enumerate(body_paragraphs):
            if i == 0 and not title:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = para
            p.font.name = FONT_SANS
            p.font.size = Pt(font_size)
            p.font.color.rgb = SLATE
            p.space_after = Pt(4)
            
        return card

    # Helper to add a structured text-block card matching the user's aesthetic
    def add_text_block_card(slide, left, top, width, height, sub_label, bold_title, body_text, accent_color=RUST, color=None):
        card_fill = color if isinstance(color, RGBColor) else CARD_BG
        
        # 1. Main Card Shape (Light Alabaster Cream rounded rectangle)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = card_fill
        card.line.color.rgb = MUTED_GRAY
        card.line.width = Pt(0.75)
        
        # 2. Top Accent Line (Solid accent color bar)
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.08))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_color
        accent_bar.line.fill.background()
        
        # 3. Inside Text Frame
        txBox = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.3), width - Inches(0.5), height - Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        # 3a. Sub-label: uppercase, small, spaced, accent color
        p_sub = tf.paragraphs[0]
        p_sub.text = sub_label.upper()
        p_sub.font.name = FONT_SANS
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = accent_color
        p_sub.space_after = Pt(6)
        
        # 3b. Bold Title
        p_title = tf.add_paragraph()
        p_title.text = bold_title
        p_title.font.name = FONT_SERIF
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = SLATE
        p_title.space_after = Pt(10)
        
        # 3c. Body Text
        p_body = tf.add_paragraph()
        p_body.text = body_text
        p_body.font.name = FONT_SANS
        p_body.font.size = Pt(14)
        p_body.font.color.rgb = SLATE
        p_body.space_after = Pt(0)
        
        return card


    # Use blank layouts
    blank_slide_layout = prs.slide_layouts[6]
    TOTAL_SLIDES = 25
    
    # ----------------------------------------------------------    # SLIDE 1: Title Slide (Cream Background with Serif/Accent Typography)
    # -------------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide1)
    
    # Title text box
    txBox = slide1.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.733), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "" # Clear default
    p.alignment = PP_ALIGN.LEFT
    
    # Run 1: Normal title text
    run1 = p.add_run()
    run1.text = "Development of a "
    run1.font.name = FONT_SERIF
    run1.font.size = Pt(38)
    run1.font.bold = True
    run1.font.color.rgb = SLATE
    
    # Run 2: Emphasized disease diagnosis (italics, rust)
    run2 = p.add_run()
    run2.text = "Heart Disease Diagnosis"
    run2.font.name = FONT_SERIF
    run2.font.size = Pt(38)
    run2.font.bold = True
    run2.font.italic = True
    run2.font.color.rgb = RUST
    
    # Run 3: Remainder of title
    run3 = p.add_run()
    run3.text = " Technique\nUsing ECG (Electrocardiogram) Signals"
    run3.font.name = FONT_SERIF
    run3.font.size = Pt(38)
    run3.font.bold = True
    run3.font.color.rgb = SLATE
    
    # Thin divider under title
    und = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.9), Inches(11.733), Inches(0.02))
    und.fill.solid()
    und.fill.fore_color.rgb = SLATE
    und.line.fill.background()
    
    # Candidate details (Left column)
    txBox_cand = slide1.shapes.add_textbox(Inches(0.8), Inches(4.3), Inches(5.0), Inches(2.2))
    tf_cand = txBox_cand.text_frame
    tf_cand.word_wrap = True
    tf_cand.margin_left = tf_cand.margin_right = tf_cand.margin_top = tf_cand.margin_bottom = 0
    
    p_c_label = tf_cand.paragraphs[0]
    p_c_label.text = "CANDIDATE"
    p_c_label.font.name = FONT_SANS
    p_c_label.font.size = Pt(11)
    p_c_label.font.bold = True
    p_c_label.font.color.rgb = MUTED_GRAY
    p_c_label.space_after = Pt(4)
    
    p_c_val = tf_cand.add_paragraph()
    p_c_val.text = "Askar Anafin"
    p_c_val.font.name = FONT_SERIF
    p_c_val.font.size = Pt(22)
    p_c_val.font.bold = True
    p_c_val.font.color.rgb = SLATE
    p_c_val.space_after = Pt(6)
    
    p_c_sub = tf_cand.add_paragraph()
    p_c_sub.text = "Applied Data Analytics (7M06103)\nSchool of Artificial Intelligence & Data Science\nAstana IT University"
    p_c_sub.font.name = FONT_SANS
    p_c_sub.font.size = Pt(13)
    p_c_sub.font.color.rgb = SLATE
    p_c_sub.space_after = Pt(4)
    
    # Supervisor details (Right column)
    txBox_sup = slide1.shapes.add_textbox(Inches(6.8), Inches(4.3), Inches(5.0), Inches(2.2))
    tf_sup = txBox_sup.text_frame
    tf_sup.word_wrap = True
    tf_sup.margin_left = tf_sup.margin_right = tf_sup.margin_top = tf_sup.margin_bottom = 0
    
    p_s_label = tf_sup.paragraphs[0]
    p_s_label.text = "SCIENTIFIC SUPERVISOR"
    p_s_label.font.name = FONT_SANS
    p_s_label.font.size = Pt(11)
    p_s_label.font.bold = True
    p_s_label.font.color.rgb = MUTED_GRAY
    p_s_label.space_after = Pt(4)
    
    p_s_val = tf_sup.add_paragraph()
    p_s_val.text = "Minsoo Hahn, PhD"
    p_s_val.font.name = FONT_SERIF
    p_s_val.font.size = Pt(22)
    p_s_val.font.bold = True
    p_s_val.font.color.rgb = SLATE
    p_s_val.space_after = Pt(6)
    
    # AITU logo (if exists, placed on the bottom right)
    logo_path = r"figures\aitu_logo.png"
    if os.path.exists(logo_path):
        slide1.shapes.add_picture(logo_path, Inches(10.5), Inches(4.3), width=Inches(2.0))


    # -------------------------------------------------------------
    # SLIDE 3: Relevance of Thesis Topic
    # -------------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide3)
    add_slide_header(slide3, "Relevance of Thesis Topic", 2, TOTAL_SLIDES)
    
    relevance_points = [
        "Global Mortality: CVDs are the leading cause of death worldwide, causing 17.9M deaths/year.",
        "Clinical Bottleneck: Manual ECG interpretation is slow, subjective, and prone to fatigue.",
        "High Variability: Inter-observer diagnostic error exceeds 20% in emergency departments.",
        "Fragile Handcrafted Features: Traditional ML relies on noise-sensitive peak detection.",
        "Spatial Gap: Standard 1D CNNs ignore the 3D anatomical geometry of electrode placement."
    ]
    add_textbox(slide3, Inches(0.8), Inches(1.6), Inches(7.5), Inches(5.2), relevance_points, size=17, bold_first_n_words=1)
    
    card_body = [
        "Develop deep learning systems that model local wave morphology and spatial lead projections for real-time, explainable triage."
    ]
    add_card(slide3, Inches(8.8), Inches(1.8), Inches(3.8), Inches(4.5), "The Research Need:", card_body, border_color=BLACK, font_size=15)

    # -------------------------------------------------------------
    # SLIDE 4: Research Aim
    # -------------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide4)
    add_slide_header(slide4, "Research Aim", 3, TOTAL_SLIDES)
    
    # Large block for Research Aim in ONE complete sentence (Plain Border box)
    aim_box = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.0))
    aim_box.fill.solid()
    aim_box.fill.fore_color.rgb = LIGHT_GRAY
    aim_box.line.color.rgb = BLACK
    aim_box.line.width = Pt(2.0)
    
    tf_aim = aim_box.text_frame
    tf_aim.word_wrap = True
    p_aim_title = tf_aim.paragraphs[0]
    p_aim_title.text = "FORMAL STATEMENT OF THE RESEARCH AIM:"
    p_aim_title.font.name = FONT_TITLE
    p_aim_title.font.size = Pt(16)
    p_aim_title.font.bold = True
    p_aim_title.font.color.rgb = BLACK
    p_aim_title.space_after = Pt(12)
    
    p_aim_text = tf_aim.add_paragraph()
    p_aim_text.text = "The main aim of this research is to develop and evaluate deep learning methods for automated multi-lead electrocardiogram classification to provide clinical decision support and patient triage."
    p_aim_text.font.name = FONT_BODY
    p_aim_text.font.size = Pt(22)
    p_aim_text.font.bold = True
    p_aim_text.font.color.rgb = BLACK
    p_aim_text.alignment = PP_ALIGN.CENTER
    
    info_points = [
        "Clinical Impact: Targets reducing the 20% emergency room observer error rate.",
        "Model Scope: Focuses on standard 12-lead digital ECG time-series signals.",
        "Deployment: Lightweight web application delivering inferences in milliseconds."
    ]
    add_textbox(slide4, Inches(1.0), Inches(4.7), Inches(11.333), Inches(2.0), info_points, size=18, bold_first_n_words=1, bullet=True)

    # -------------------------------------------------------------
    # SLIDE 5: Objectives
    # -------------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide5)
    add_slide_header(slide5, "Objectives", 4, TOTAL_SLIDES)
    
    # Helper for objectives cards (simple border boxes with white background)
    def add_obj_card(slide, left, top, width, height, number_str, text_str):
        # Draw a rounded rectangle
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = MUTED_GRAY
        card.line.width = Pt(0.75)
        
        # Add textbox inside
        txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        # Number prefix in Rust and bold
        run_num = p.add_run()
        run_num.text = number_str + "  "
        run_num.font.name = FONT_SANS
        run_num.font.bold = True
        run_num.font.size = Pt(14)
        run_num.font.color.rgb = RUST
        
        # Text in Slate
        run_text = p.add_run()
        run_text.text = text_str
        run_text.font.name = FONT_SANS
        run_text.font.size = Pt(12)
        run_text.font.color.rgb = SLATE
        p.alignment = PP_ALIGN.LEFT

    # Left Column (Objectives 1-4)
    add_obj_card(slide5, Inches(0.8), Inches(1.5), Inches(5.6), Inches(1.1), "①", "Build a zero-phase preprocessing pipeline to filter baseline wander and high-frequency noise.")
    add_obj_card(slide5, Inches(0.8), Inches(2.72), Inches(5.6), Inches(1.1), "②", "Develop a 1D ResNet-18 model as a baseline for wave morphology extraction.")
    add_obj_card(slide5, Inches(0.8), Inches(3.94), Inches(5.6), Inches(1.1), "③", "Design ST-ReGE GNN with a learnable adjacency matrix to capture spatial lead projections.")
    add_obj_card(slide5, Inches(0.8), Inches(5.16), Inches(5.6), Inches(1.1), "④", "Develop a 1D Vision Transformer model to capture long-range rhythm dependencies.")

    # Right Column (Objectives 5-7)
    add_obj_card(slide5, Inches(6.8), Inches(1.5), Inches(5.7), Inches(1.4), "⑤", "Formulate a Spatio-Temporal Hybrid combining convolutional morphology with transformer rhythm.")
    add_obj_card(slide5, Inches(6.8), Inches(3.1), Inches(5.7), Inches(1.4), "⑥", "Calibrate boundaries via post-hoc validation search and soft-voting ensembling.")
    add_obj_card(slide5, Inches(6.8), Inches(4.7), Inches(5.7), Inches(1.4), "⑦", "Build a real-time FastAPI web dashboard for triage upload and wave visualization.")

    # -------------------------------------------------------------
    # SLIDE 6: Literature Review
    # -------------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide6)
    add_slide_header(slide6, "Literature Review", 5, TOTAL_SLIDES)
    
    # Create Table
    rows, cols = 8, 5
    t_shape = slide6.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(4.5))
    table = t_shape.table
    
    # Set widths
    table.columns[0].width = Inches(2.8)
    table.columns[1].width = Inches(2.4)
    table.columns[2].width = Inches(2.1)
    table.columns[3].width = Inches(2.8)
    table.columns[4].width = Inches(2.233)
    
    headers = ["Publication Authors & Year", "Data Acquisition Method", "Dataset Size", "Model(s) Used", "Accuracy / Performance"]
    lit_data = [
        ["Kiranyaz et al. (2016)", "Bipolar ECG (MIT-BIH)", "48 records", "1D CNN (patient-specific)", "98.9% (beat-wise acc)"],
        ["Rajpurkar et al. (2017)", "Single-lead (Zio Patch)", "91,232 records", "34-layer deep 1D CNN", "F1-score: 0.837"],
        ["Ribeiro et al. (2020)", "12-lead (Telehealth Minas)", "2,322,513 records", "1D ResNet-based CNN", "F1-score > 0.80"],
        ["Zhang et al. (2021)", "12-lead (CPSC 2018)", "6,877 records", "ECG-GNN (fixed spatial)", "F1-score: 0.812"],
        ["Vaid et al. (2023)", "12-lead (Mount Sinai)", "1,077,137 records", "1D Vision Transformer", "AUROC: 0.940"],
        ["Strodthoff et al. (2021)", "12-lead (PTB-XL benchmark)", "21,837 records", "1D ResNet-18 benchmark", "Macro F1: 0.705–0.720"],
        ["This Thesis (June 2026)", "12-lead (PTB-XL, filtered)", "21,388 records", "ST-ReGE GNN (learnable) + Ensemble & Calibration", "Macro F1: 0.7506 / AUROC: 0.9344"]
    ]
    
    # Headers formatting
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLACK
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_TITLE
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
    # Data formatting
    for r_idx, r_data in enumerate(lit_data):
        for c_idx, val in enumerate(r_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            # Highlight this thesis using thin gray shading
            if r_idx == 6:
                cell.fill.fore_color.rgb = LIGHT_GRAY
            elif r_idx % 2 == 0:
                cell.fill.fore_color.rgb = WHITE
            else:
                cell.fill.fore_color.rgb = LIGHT_GRAY
                
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_BODY
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = BLACK
                if r_idx == 6:
                    p.font.bold = True
                    
    # Bottom note
    note_box = slide6.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(12.333), Inches(0.8))
    tf_n = note_box.text_frame
    p_n = tf_n.paragraphs[0]
    p_n.text = "Research Choice: ST-ReGE GNN and ensembling model lead geometry and rhythms, outperforming static convolutional baselines."
    p_n.font.name = FONT_BODY
    p_n.font.size = Pt(12)
    p_n.font.italic = True
    p_n.font.color.rgb = BLACK

    # -------------------------------------------------------------
    # SLIDE 7: Novelty of the Research (Redesigned: 3 Columns + Callout)
    # -------------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide7)
    add_slide_header(slide7, "Novelty of the Research", 6, TOTAL_SLIDES)
    
    # Card 1: Method (Rust accent)
    add_text_block_card(slide7, Inches(0.8), Inches(1.6), Inches(3.6), Inches(4.0), 
                        "01  METHODOLOGY", 
                        "Dynamic Lead Graph", 
                        "Models standard leads as graph nodes rather than independent stacked channels. A learnable, sigmoid-activated adjacency matrix dynamically discovers anatomical pathways and electrical correlations during training.", 
                        accent_color=RUST)
             
    # Card 2: Evaluation (Teal accent)
    add_text_block_card(slide7, Inches(4.8), Inches(1.6), Inches(3.6), Inches(4.0), 
                        "02  EVALUATION", 
                        "Multi-Paradigm Comparison", 
                        "Provides a fair, standardized benchmark evaluating local convolutional morphology, lead-spatial propagation, and global attention models on the same filtered dataset to isolate true architectural benefits.", 
                        accent_color=TEAL)
             
    # Card 3: Optimization (Peach/Orange accent)
    add_text_block_card(slide7, Inches(8.8), Inches(1.6), Inches(3.6), Inches(4.0), 
                        "03  CALIBRATION", 
                        "Decision Boundary Tuning", 
                        "Applies validation-set threshold grid search to calibrate decision boundaries for minority classes. This handles severe label imbalances and optimizes clinical recall in seconds without retraining.", 
                        accent_color=PEACH)

             
    # Bottom Callout Box (Light background with left vertical bar)
    c_box = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.8), Inches(11.733), Inches(0.85))
    c_box.fill.solid()
    c_box.fill.fore_color.rgb = CARD_BG
    c_box.line.color.rgb = MUTED_GRAY
    c_box.line.width = Pt(0.75)
    
    bar_c = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.8), Inches(0.08), Inches(0.85))
    bar_c.fill.solid()
    bar_c.fill.fore_color.rgb = RUST
    bar_c.line.fill.background()
    
    txBox_c = slide7.shapes.add_textbox(Inches(1.0), Inches(5.9), Inches(11.333), Inches(0.65))
    tf_c = txBox_c.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_right = tf_c.margin_top = tf_c.margin_bottom = 0
    p_c = tf_c.paragraphs[0]
    p_c.text = "Addresses three key gaps: physical lead layout modeling, model comparison fairness, and screening threshold calibration."
    p_c.font.name = FONT_SANS
    p_c.font.size = Pt(14)
    p_c.font.italic = True
    p_c.font.color.rgb = SLATE

    # -------------------------------------------------------------
    # SLIDE 8: Problem Statement: Input/Output
    # -------------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide8)
    add_slide_header(slide8, "Problem Statement: Input/Output Tensors", 7, TOTAL_SLIDES)
    
    input_desc = [
        "Digital Signal Tensor: X in R^(B x 12 x 5000)",
        "B: Batch size of records.",
        "12 Leads: Standard chest and limb electrode recordings.",
        "5000 Steps: 10 seconds of signal sampled at 500 Hz."
    ]
    
    output_desc = [
        "Binary Multi-Label Vector: y in {0, 1}^(B x 5)",
        "5 Superclasses: NORM, MI, STTC, CD, HYP",
        "Co-occurring diagnoses: Allowed (e.g. MI + CD co-exist)."
    ]
    
    add_card(slide8, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "Input Data (ECG Signal)", input_desc, font_size=16, color=WHITE)
    add_card(slide8, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), "Output Data (Diagnostic Labels)", output_desc, border_color=BLACK, font_size=16, color=WHITE)

    # -------------------------------------------------------------
    # SLIDE 8b: ECG Signal Data Table
    # -------------------------------------------------------------
    slide8b = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide8b)
    add_slide_header(slide8b, "ECG Signal Data Table (12 Leads)", 8, TOTAL_SLIDES)
    
    # Text explanation above table
    add_textbox(slide8b, Inches(0.5), Inches(1.2), Inches(12.333), Inches(0.5), 
                ["Sample values of ECG signal tensor X (values in mV, sampled at 500 Hz):"], size=16, bullet=False)
    
    # Table of ECG values: 6 rows (1 header + 5 time steps), 13 columns (Time + 12 leads)
    rows, cols = 6, 13
    table_shape = slide8b.shapes.add_table(rows, cols, Inches(0.5), Inches(1.8), Inches(12.333), Inches(3.2))
    table = table_shape.table
    
    # Set column widths
    table.columns[0].width = Inches(1.533) # Time Step column
    for c in range(1, 13):
        table.columns[c].width = Inches(0.9)
        
    headers = ["Time Step", "I", "II", "III", "aVR", "aVL", "aVF", "V1", "V2", "V3", "V4", "V5", "V6"]
    
    # Table values
    table_data = [
        ["0 ms (t1)", "-0.08", "0.15", "0.23", "-0.04", "-0.16", "0.19", "-0.05", "0.12", "0.35", "0.82", "0.55", "0.21"],
        ["2 ms (t2)", "-0.07", "0.18", "0.25", "-0.05", "-0.16", "0.21", "-0.04", "0.15", "0.38", "0.88", "0.58", "0.23"],
        ["4 ms (t3)", "-0.05", "0.22", "0.27", "-0.07", "-0.16", "0.24", "-0.03", "0.18", "0.42", "0.95", "0.62", "0.26"],
        ["6 ms (t4)", "-0.03", "0.26", "0.29", "-0.09", "-0.16", "0.27", "-0.02", "0.21", "0.46", "1.02", "0.66", "0.29"],
        ["8 ms (t5)", "-0.01", "0.30", "0.31", "-0.11", "-0.16", "0.30", "-0.01", "0.24", "0.50", "1.08", "0.70", "0.32"]
    ]
    
    # Header format
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLACK
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_TITLE
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
    # Data format
    for r_idx, r_data in enumerate(table_data):
        for c_idx, val in enumerate(r_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            if r_idx % 2 == 0:
                cell.fill.fore_color.rgb = WHITE
            else:
                cell.fill.fore_color.rgb = LIGHT_GRAY
                
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_BODY
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = BLACK
                
    # Add block text below
    add_card(slide8b, Inches(0.5), Inches(5.3), Inches(12.333), Inches(1.5), "Tensor Interpretation:",
             ["Each row represents a single time step separated by 2 ms.",
              "The columns represent standard 12-lead potential differences in millivolts (mV).",
              "A single 10-second ECG record contains 5000 such rows (B x 12 x 5000 tensor shape)."], font_size=13)

    # -------------------------------------------------------------
    # SLIDE 9: Collection of Input Data: Initial Dataset
    # -------------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide9)
    add_slide_header(slide9, "Initial Dataset Overview", 9, TOTAL_SLIDES)
    
    initial_desc = [
        "Database: PTB-XL corpus containing 21,799 clinical 12-lead ECG records from 18,869 patients.",
        "Demographics: 52% are male and 48% are female. Ages cover 0 to 95 years (median 62, interquartile range of 22).",
        "Patient Stratification: Fold column (strat_fold 1-10) isolates patient records to prevent leakage and preserves diagnostic label ratios across splits.",
        "Class Imbalance: Severe ~4:1 ratio (Normal vs. Hypertrophy) requires post-hoc boundary calibration to optimize minority diagnostic recall."
    ]
    add_textbox(slide9, Inches(0.8), Inches(1.6), Inches(6.8), Inches(5.0), initial_desc, size=17, bold_first_n_words=1)
    
    # Class distribution card on the right
    dist_desc = [
        "NORM (Normal): 9,528 records",
        "MI (Myocardial Infarction): 5,486 records",
        "STTC (ST/T Changes): 5,250 records",
        "CD (Conduction Disturbance): 4,907 records",
        "HYP (Hypertrophy): 2,655 records",
        "",
        "Total diagnostic statements: 71"
    ]
    add_card(slide9, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.5), "Class-wise Distribution", dist_desc, border_color=BLACK, font_size=16, color=WHITE)

    # -------------------------------------------------------------
    # SLIDE 10: Processing of Input Data: Preprocessing
    # -------------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide10)
    add_slide_header(slide10, "Zero-Phase Signal Preprocessing", 10, TOTAL_SLIDES)
    
    pre_desc = [
        "Noise Sources: Baseline wander (<0.5 Hz) from breathing and high-frequency EMG muscle noise.",
        "Bandpass Filter: 4th-order Butterworth digital filter (0.5 to 50 Hz).",
        "Zero-Phase Delay: Forward-backward filtering (filtfilt) cancels out phase shifts and prevents wave distortion.",
        "Normalization: Lead-wise Z-score scaling standardizes amplitudes across records."
    ]
    add_textbox(slide10, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), pre_desc, size=17, bold_first_n_words=1)
    
    # Accent container card for preprocessing comparison image
    add_card(slide10, Inches(6.8), Inches(1.6), Inches(5.7), Inches(4.8), "", [], accent_color=TEAL)
    
    # Insert preprocessing comparison image
    comparison_img_path = r"figures\preprocessing_comparison.png"
    if os.path.exists(comparison_img_path):
        slide10.shapes.add_picture(comparison_img_path, Inches(6.95), Inches(1.85), width=Inches(5.4), height=Inches(4.3))
    else:
        txBox_e = slide10.shapes.add_textbox(Inches(6.95), Inches(1.85), Inches(5.4), Inches(4.3))
        tf_e = txBox_e.text_frame
        tf_e.word_wrap = True
        p_e = tf_e.paragraphs[0]
        p_e.text = "[preprocessing_comparison.png placeholder]"
        p_e.font.name = FONT_SANS
        p_e.font.size = Pt(14)
        p_e.font.bold = True
        p_e.font.color.rgb = SLATE
        p_e.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 11: Cleaning of Input Data: Label Filtering
    # -------------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide11)
    add_slide_header(slide11, "Diagnostic Label Filtering", 11, TOTAL_SLIDES)
    
    cleaning_desc = [
        "Diagnostic Filtering: Retains only records with active diagnostic statement labels.",
        "All-Zeros Bias: 411 records (1.89%) have only rhythm/form codes, yielding all-zero labels.",
        "Metric Inflation: Keeping all-zero records creates optimistic bias, artificially inflating F1 and AUPRC.",
        "Final Dataset: Removing all-zero records leaves 21,388 records, matching standard benchmarks.",
        "Class Totals: Positive class counts are unaffected (removed records were completely empty)."
    ]
    add_textbox(slide11, Inches(0.8), Inches(1.6), Inches(6.8), Inches(5.0), cleaning_desc, size=17, bold_first_n_words=1)
    
    # Before/After table on the right
    c_rows, c_cols = 8, 3
    t_c_shape = slide11.shapes.add_table(c_rows, c_cols, Inches(8.0), Inches(1.8), Inches(4.5), Inches(4.0))
    tc = t_c_shape.table
    tc.columns[0].width = Inches(1.8)
    tc.columns[1].width = Inches(1.35)
    tc.columns[2].width = Inches(1.35)
    
    tc_headers = ["Category", "Raw Count", "Filtered"]
    tc_data = [
        ["Full Corpus", "21,799", "21,388"],
        ["Normal (NORM)", "9,514", "9,514"],
        ["Infarction (MI)", "5,469", "5,469"],
        ["ST/T Changes", "5,235", "5,235"],
        ["Blocks (CD)", "4,898", "4,898"],
        ["Hypertrophy", "2,649", "2,649"],
        ["All-Zero Target", "411", "0"]
    ]
    
    for col_idx, header in enumerate(tc_headers):
        cell = tc.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLACK
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_TITLE
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
    for r_idx, r_data in enumerate(tc_data):
        for c_idx, val in enumerate(r_data):
            cell = tc.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            if r_idx % 2 == 0:
                cell.fill.fore_color.rgb = WHITE
            else:
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_BODY
                p.font.size = Pt(13)
                p.font.color.rgb = BLACK
                p.alignment = PP_ALIGN.CENTER
                if r_idx == 6:
                    p.font.bold = True

    # -------------------------------------------------------------
    # SLIDE 12: Models and Methods: Baseline Models: CNN & ViT
    # -------------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide12)
    add_slide_header(slide12, "Baseline Models: CNN & ViT", 12, TOTAL_SLIDES)
    
    cnn_desc = [
        "Temporal Morphology: Convolutions capture wave shapes along the temporal axis.",
        "Inductive Bias: Weight-sharing and translation invariance make CNNs highly data-efficient.",
        "Architecture: Stacks four residual stages with shortcut projections to prevent gradient degradation."
    ]
    
    vit_desc = [
        "Sequential Rhythm: Models long-range temporal dependencies and cardiac rhythms using self-attention.",
        "Tokenization: Splits temporal sequence into 100 non-overlapping patches, prepending a [CLS] token.",
        "Constraint: Lacks spatial and temporal locality biases, making it data-hungry and prone to overfitting."
    ]
    
    add_card(slide12, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "1D Adapted ResNet-18 (CNN)", cnn_desc, font_size=15, color=WHITE, border_color=BLACK)
    add_card(slide12, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), "1D Vision Transformer (ViT-1D)", vit_desc, border_color=BLACK, font_size=15, color=WHITE)

    # -------------------------------------------------------------
    # SLIDE 13: Models and Methods: ST-ReGE GNN
    # -------------------------------------------------------------
    slide13 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide13)
    add_slide_header(slide13, "Spatio-Temporal ReGE GNN", 13, TOTAL_SLIDES)
    
    gnn_left = [
        "Lead-Spatial Graph: Models 12 physical leads as graph nodes to capture spatial correlations.",
        "Shared Backbone: Independent 1D CNN processes each lead to extract 256-D morphological features.",
        "Learnable Adjacency: Trainable A = Sigmoid(W_adj) in R^12x12 dynamically discovers functional electrical pathways.",
        "Message Passing: Two GCN layers propagate features over the discovered pathways."
    ]
    add_textbox(slide13, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), gnn_left, size=17, bold_first_n_words=1)
    
    gnn_right = [
        "Tuning & Optimization Approach:",
        "    • Hyperparameters: Batch size = 32, patience = 10 epochs",
        "    • Optimizer: Adam (LR = 10^-3, Weight Decay = 10^-4)",
        "    • Scheduler: Cosine Annealing learning rate schedule",
        "",
        "GCN Feature Propagation Formula:",
        "    H(l+1) = BatchNorm( ReLU( A * H(l) * W(l) ) ) + H(l)",
        "    • A: Learnable Adjacency, W(l): Trainable weights",
        "",
        "Theoretical Concept:",
        "    • Learnable Adjacency: Discovers functional connections directly from data, rather than assuming physical lead geometry is rigid.",
        "    • Simple Analogy: A social network where leads learn to talk directly to anyone whose signals correlate, rather than only to neighbors."
    ]
    add_card(slide13, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.5), "ST-ReGE Concept & Tuning", gnn_right, border_color=BLACK, font_size=13, color=WHITE)

    # -------------------------------------------------------------
    # SLIDE 14: Models and Methods: Spatio-Temporal Hybrid Model
    # -------------------------------------------------------------
    slide14 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide14)
    add_slide_header(slide14, "Spatio-Temporal Hybrid Model", 14, TOTAL_SLIDES)
    
    hybrid_left = [
        "CNN Front-End: 3 ResNet stages compress raw signals to 313 steps to extract local QRS morphology.",
        "Squeeze-and-Excitation: SE channel attention blocks recalibrate lead importance before temporal modeling.",
        "Transformer Back-End: 3-layer Encoder (8 heads, d_model = 256) models global rhythm context."
    ]
    add_textbox(slide14, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), hybrid_left, size=17, bold_first_n_words=1)
    
    hybrid_right = [
        "Regularization & Tuning Approach:",
        "    • Hyperparameters: Batch size = 32, patience = 10 epochs",
        "    • Classifier: 2-layer MLP with 0.2 dropout",
        "    • Optimizer: Adam (LR = 10^-3, Weight Decay = 10^-4)",
        "    • Scheduler: Cosine Annealing learning rate schedule",
        "",
        "Spatio-Temporal Integration Synergy:",
        "    • Combines local morphology bias of CNN with global temporal rhythm checks of Self-Attention.",
        "    • Calibrates lead attention (SE blocks) prior to sequence modeling.",
        "",
        "Analogy (Spellchecker vs. Speed-Reader):",
        "    • CNN front-end acts as a spellchecker to verify individual beat shapes.",
        "    • Transformer back-end speed-reads overall rhythm patterns."
    ]
    add_card(slide14, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.5), "Hybrid Concept & Tuning", hybrid_right, border_color=BLACK, font_size=13, color=WHITE)

    # -------------------------------------------------------------
    # SLIDE 14b: Core Outcomes of ST-ReGE & Hybrid
    # -------------------------------------------------------------
    slide14b = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide14b)
    add_slide_header(slide14b, "Core Outcomes: ST-ReGE & Hybrid", 15, TOTAL_SLIDES)
    
    gnn_outcomes = [
        "How it Works: Models 12 leads as graph nodes and dynamically learns physical adjacency A = Sigmoid(W_adj) combined with 2 GCN feature propagation layers.",
        "What it Solves: The spatial lead propagation gap of traditional 1D CNNs, discovering functional electrical pathways across different cardiac walls.",
        "Theoretical Value: Captures joint electrical changes across leads (e.g. inferior leads II, III, aVF) instead of treating leads as independent channels.",
        "Limitations: Assumes lead geometry fits a learnable static graph; increased model complexity compared to standard baseline convolutions."
    ]
    
    hybrid_outcomes = [
        "How it Works: Stacks a 3-stage ResNet morphological front-end, SE channel lead attention block, and a 3-layer Transformer Encoder rhythm back-end.",
        "What it Solves: The data-hunger and overfitting constraints of pure 1D Vision Transformers on medium-sized database splits.",
        "Theoretical Value: Blends the local temporal inductive bias of CNNs (beat shape details) with the global rhythm sequence attention of Transformers.",
        "Limitations: Higher computational memory and parameter counts; multi-head self-attention adds latency compared to simple convolutions."
    ]
    
    add_card(slide14b, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), "ST-ReGE GNN Outcomes", gnn_outcomes, font_size=13, color=WHITE, border_color=BLACK, accent_color=RUST)
    add_card(slide14b, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), "Spatio-Temporal Hybrid Outcomes", hybrid_outcomes, font_size=13, color=WHITE, border_color=BLACK, accent_color=TEAL)

    # -------------------------------------------------------------
    # SLIDE 15: Ensemble & Boundary Calibration (Redesigned: 3 Columns + Callout)
    # -------------------------------------------------------------
    slide15 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide15)
    add_slide_header(slide15, "Ensemble & Boundary Calibration", 16, TOTAL_SLIDES)
    
    # Card 1: Soft Ensembling (Rust accent)
    add_text_block_card(slide15, Inches(0.8), Inches(1.6), Inches(3.6), Inches(4.0), 
                        "SOFT ENSEMBLING", 
                        "Joint Probability Voting", 
                        "Blends CNN morphology, GNN spatial routing, and Hybrid attention predictions by averaging their output probabilities. This stabilizes final classifications and mitigates model-specific bias.", 
                        accent_color=RUST)
             
    # Card 2: Threshold Calibration (Teal accent)
    add_text_block_card(slide15, Inches(4.8), Inches(1.6), Inches(3.6), Inches(4.0), 
                        "THRESHOLD CALIBRATION", 
                        "Validation-Set Grid Search", 
                        "Replaces default 0.5 boundaries with class-specific thresholds: [NORM: 0.48, MI: 0.42, STTC: 0.34, CD: 0.47, HYP: 0.39]. Shifting boundaries down boosts recall on minority classes.", 
                        accent_color=TEAL)
             
    # Card 3: Clinical Utility (Peach/Orange accent)
    add_text_block_card(slide15, Inches(8.8), Inches(1.6), Inches(3.6), Inches(4.0), 
                        "CLINICAL UTILITY", 
                        "On-the-Fly Adaptability", 
                        "Allows clinicians to tune screening sensitivity dynamically to match ICU triage requirements. Grid search runs in seconds on CPU without retraining network weights.", 
                        accent_color=PEACH)
              
    # Bottom Callout Box (Light background with left vertical bar)
    c_box = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.8), Inches(11.733), Inches(0.85))
    c_box.fill.solid()
    c_box.fill.fore_color.rgb = CARD_BG
    c_box.line.color.rgb = MUTED_GRAY
    c_box.line.width = Pt(0.75)
    
    bar_c = slide15.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(5.8), Inches(0.08), Inches(0.85))
    bar_c.fill.solid()
    bar_c.fill.fore_color.rgb = RUST
    bar_c.line.fill.background()
    
    txBox_c = slide15.shapes.add_textbox(Inches(1.0), Inches(5.9), Inches(11.333), Inches(0.65))
    tf_c = txBox_c.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_right = tf_c.margin_top = tf_c.margin_bottom = 0
    p_c = tf_c.paragraphs[0]
    p_c.text = "Why Calibration over Loss-Weighting? It preserves probability rank-ordering capability (AUPRC) and runs in seconds without altering model gradients."
    p_c.font.name = FONT_SANS
    p_c.font.size = Pt(14)
    p_c.font.italic = True
    p_c.font.color.rgb = SLATE

    # -------------------------------------------------------------
    # SLIDE 15b: Why Multi-Label Evaluation Metrics Matter
    # -------------------------------------------------------------
    slide15b = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide15b)
    add_slide_header(slide15b, "Why Multi-Label Evaluation Metrics Matter", 17, TOTAL_SLIDES)
    
    metrics_left = [
        "Class Imbalance Handling: ECG datasets like PTB-XL are highly imbalanced (e.g. NORM vs. HYP is ~4:1). Accuracy is misleading; predicting all negative for HYP gets 88% accuracy but is clinically useless.",
        "Macro F1-Score: Calculates the balanced harmonic mean of Precision and Recall for each class independently, then averages them. Treats rare but critical pathologies (such as Hypertrophy) with equal weight.",
        "AUROC (Area Under ROC): Evaluates threshold-independent discriminative capacity; the probability that the model ranks a random positive case higher than a random negative.",
        "AUPRC (Area Under PR Curve): Focuses heavily on the positive class (pathology), preventing false positive/negative inflation from large true negative pools."
    ]
    add_textbox(slide15b, Inches(0.8), Inches(1.6), Inches(6.2), Inches(5.0), metrics_left, size=16, bold_first_n_words=1)
    
    metrics_right = [
        "Clinical Triage Implications:",
        "    • High Recall (Sensitivity): Crucial for triage screening to ensure no active cardiac pathology (e.g., MI) is missed (False Negative = 0).",
        "    • High Precision: Required to prevent false alarms that overwhelm emergency room staff (False Positive = 0).",
        "",
        "Alignment with Clinical Utility:",
        "    • PhysioNet Challenge Score uses a clinical weight matrix to penalize high-risk misclassifications (e.g., calling MI Normal) more severely than low-risk ones.",
        "    • Threshold grid search allows clinicians to adjust the model's operating point on-the-fly to balance recall and precision."
    ]
    add_card(slide15b, Inches(7.4), Inches(1.8), Inches(5.1), Inches(4.5), "Clinical Metric Alignment", metrics_right, border_color=BLACK, font_size=13, color=WHITE)

    # SLIDE 16: Results: Quantitative Benchmarks
    # -------------------------------------------------------------
    slide16 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide16)
    add_slide_header(slide16, "Quantitative Model Benchmarks", 18, TOTAL_SLIDES)
    
    # Benchmark Table
    rows, cols = 6, 7
    table_shape = slide16.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(4.5))
    table = table_shape.table
    
    # Set widths
    table.columns[0].width = Inches(2.7)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(1.5)
    table.columns[3].width = Inches(1.5)
    table.columns[4].width = Inches(1.5)
    table.columns[5].width = Inches(1.7)
    table.columns[6].width = Inches(1.933)
    
    headers = ["Metric", "CNN (ResNet-18)", "GNN (ST-ReGE)", "ViT-1D", "S-T Hybrid", "Ensemble (Default)", "Ensemble (Opt.)"]
    bench_data = [
        ["Macro F1-Score", "0.7285", "0.7444", "0.6820", "0.7226", "0.7435", "0.7506 (+0.71%)"],
        ["Macro AUROC", "0.9266", "0.9269", "0.8939", "0.9221", "0.9344", "0.9344"],
        ["Macro AUPRC", "0.8125", "0.8218", "0.7484", "0.8064", "0.8354", "0.8354"],
        ["Exact Match Acc", "63.58%", "64.42%", "59.51%", "62.41%", "65.59%", "64.24%"],
        ["Training Time", "22.5 min", "23.9 min", "19.5 min", "16.7 min", "N/A", "N/A"]
    ]
    
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLACK
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_TITLE
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
    for r_idx, r_data in enumerate(bench_data):
        for c_idx, val in enumerate(r_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            if c_idx == 6 and r_idx in [0, 1, 2]:
                cell.fill.fore_color.rgb = LIGHT_GRAY
            elif r_idx % 2 == 0:
                cell.fill.fore_color.rgb = WHITE
            else:
                cell.fill.fore_color.rgb = LIGHT_GRAY
                
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_BODY
                p.font.size = Pt(13)
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = BLACK
                if c_idx == 6 and r_idx in [0, 1, 2]:
                    p.font.bold = True
                elif c_idx == 2 and r_idx == 3: # GNN wins exact match
                    p.font.bold = True
                else:
                    p.font.bold = False
 
    # Bottom note
    bottom_box = slide16.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(12.333), Inches(0.8))
    tf_b = bottom_box.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "Result Summary: ST-ReGE GNN is the top single model. The calibrated soft ensemble achieves the highest overall score (0.7506 F1)."
    p_b.font.name = FONT_BODY
    p_b.font.size = Pt(13)
    p_b.font.italic = True
    p_b.font.color.rgb = SLATE

    # -------------------------------------------------------------
    # SLIDE 16b: Results: Model Interpretability (1D Grad-CAM)
    # -------------------------------------------------------------
    slide16b = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide16b)
    add_slide_header(slide16b, "Model Interpretability (1D Grad-CAM)", 19, TOTAL_SLIDES)
    
    gradcam_left = [
        "1D Grad-CAM: Computes gradients of target class scores relative to activation maps of the last convolutional layer.",
        "MI Localization: For Myocardial Infarction predictions, the peak activation is localized directly on the ST-segment elevation in Lead V2.",
        "Clinical Integration: Heatmap peaks focus on diagnostic wave morphology rather than noise, confirming the model acts as an interpretable diagnostic assistant."
    ]
    add_textbox(slide16b, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), gradcam_left, size=17, bold_first_n_words=1)
    
    # Accent container card for Grad-CAM plot
    add_card(slide16b, Inches(7.2), Inches(1.6), Inches(5.3), Inches(4.8), "", [], accent_color=TEAL)
    
    # Insert Grad-CAM image
    gradcam_img_path = r"figures\gradcam_mi.png"
    if os.path.exists(gradcam_img_path):
        slide16b.shapes.add_picture(gradcam_img_path, Inches(7.35), Inches(1.85), width=Inches(5.0), height=Inches(4.3))
    else:
        txBox_e = slide16b.shapes.add_textbox(Inches(7.35), Inches(1.85), Inches(5.0), Inches(4.3))
        tf_e = txBox_e.text_frame
        tf_e.word_wrap = True
        p_e = tf_e.paragraphs[0]
        p_e.text = "[gradcam_mi.png placeholder]"
        p_e.font.name = FONT_SANS
        p_e.font.size = Pt(14)
        p_e.font.bold = True
        p_e.font.color.rgb = SLATE
        p_e.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 16c: Results: Clinical Error Analysis (Confusion Matrix)
    # -------------------------------------------------------------
    slide16c = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide16c)
    add_slide_header(slide16c, "Clinical Error Analysis (Confusion Matrix)", 20, TOTAL_SLIDES)
    
    confusion_left = [
        "Healthy Patient Filtering: Normal Sinus Rhythms (NORM) are correctly filtered with a high 91% sensitivity.",
        "HYP vs. STTC Mimicry: 15% of Hypertrophy (HYP) cases are misclassified as ST/T changes (STTC).",
        "Ventricular Strain Justification: Left Ventricular Hypertrophy (LVH) causes mechanical strain on the heart wall, resulting in lateral ST-segment depression and T-wave inversion. This mimics primary ischemic ST/T changes, causing logical overlap."
    ]
    add_textbox(slide16c, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), confusion_left, size=17, bold_first_n_words=1)
    
    # Accent container card for Confusion Matrix plot
    add_card(slide16c, Inches(7.2), Inches(1.6), Inches(5.3), Inches(4.8), "", [], accent_color=RUST)
    
    # Insert Confusion Matrix image
    confusion_img_path = r"figures\confusion_matrix.png"
    if os.path.exists(confusion_img_path):
        slide16c.shapes.add_picture(confusion_img_path, Inches(7.35), Inches(1.85), width=Inches(5.0), height=Inches(4.3))
    else:
        txBox_e = slide16c.shapes.add_textbox(Inches(7.35), Inches(1.85), Inches(5.0), Inches(4.3))
        tf_e = txBox_e.text_frame
        tf_e.word_wrap = True
        p_e = tf_e.paragraphs[0]
        p_e.text = "[confusion_matrix.png placeholder]"
        p_e.font.name = FONT_SANS
        p_e.font.size = Pt(14)
        p_e.font.bold = True
        p_e.font.color.rgb = SLATE
        p_e.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 17: Results: FastAPI Web Dashboard
    # -------------------------------------------------------------
    slide17 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide17)
    add_slide_header(slide17, "FastAPI Diagnostic Triage Dashboard", 21, TOTAL_SLIDES)
    
    t_desc = [
        "Async Backend: FastAPI server utilizing Pydantic for validation and fast inference.",
        "File Ingestion: Validates and accepts digital ECG input files in JSON or CSV formats.",
        "Triage Dashboard: Clean light-theme interface with slate gray panels to prevent eye strain.",
        "Waveform Plots: Plots lead waveforms dynamically on pink ECG thermal grid layouts.",
        "Triage Output: Displays ensemble class probabilities and flags critical cases."
    ]
    add_textbox(slide17, Inches(0.8), Inches(1.6), Inches(6.0), Inches(5.0), t_desc, size=17, bold_first_n_words=1)
    
    # Accent container card for Web app screenshot
    add_card(slide17, Inches(7.2), Inches(1.6), Inches(5.3), Inches(4.8), "", [], accent_color=PEACH)
    
    # Insert Web app screenshot
    web_img_path = r"figures\test_webapp.png"
    if os.path.exists(web_img_path):
        slide17.shapes.add_picture(web_img_path, Inches(7.35), Inches(1.85), width=Inches(5.0), height=Inches(4.3))
    else:
        txBox_e = slide17.shapes.add_textbox(Inches(7.35), Inches(1.85), Inches(5.0), Inches(4.3))
        tf_e = txBox_e.text_frame
        tf_e.word_wrap = True
        p_e = tf_e.paragraphs[0]
        p_e.text = "[test_webapp.png placeholder]"
        p_e.font.name = FONT_SANS
        p_e.font.size = Pt(14)
        p_e.font.bold = True
        p_e.font.color.rgb = SLATE
        p_e.alignment = PP_ALIGN.CENTER
 
    # -------------------------------------------------------------
    # SLIDE 18: Scholarly Publications Supporting Thesis
    # -------------------------------------------------------------
    slide18 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide18)
    add_slide_header(slide18, "Scholarly Publication", 22, TOTAL_SLIDES)
    
    # Accent container card for publication citation on the left
    add_card(slide18, Inches(0.8), Inches(1.6), Inches(5.8), Inches(4.8), "", [], accent_color=RUST)
    
    # Text box for citation inside the card
    txBox_cit = slide18.shapes.add_textbox(Inches(1.1), Inches(2.2), Inches(5.2), Inches(3.6))
    tf_cit = txBox_cit.text_frame
    tf_cit.word_wrap = True
    tf_cit.margin_left = tf_cit.margin_right = tf_cit.margin_top = tf_cit.margin_bottom = 0
    
    # 1. Journal Header
    p_jou = tf_cit.paragraphs[0]
    p_jou.text = "SCHOLARLY CITATION"
    p_jou.font.name = FONT_SANS
    p_jou.font.size = Pt(12)
    p_jou.font.bold = True
    p_jou.font.color.rgb = RUST
    p_jou.space_after = Pt(18)
    
    # 2. Paper Title
    p_tit = tf_cit.add_paragraph()
    p_tit.text = "“Development of a Heart Disease Diagnosis Technique Using ECG Signals”"
    p_tit.font.name = FONT_SERIF
    p_tit.font.size = Pt(21)
    p_tit.font.bold = True
    p_tit.font.color.rgb = SLATE
    p_tit.space_after = Pt(18)
    
    # 3. Authors & Date
    p_auth = tf_cit.add_paragraph()
    p_auth.text = "Askar Anafin, Minsoo Hahn"
    p_auth.font.name = FONT_SANS
    p_auth.font.size = Pt(14)
    p_auth.font.color.rgb = SLATE
    p_auth.space_after = Pt(10)
    
    # 4. Journal Info
    p_status = tf_cit.add_paragraph()
    p_status.text = "Journal: Central Asian Journal (CAJ), 2026"
    p_status.font.name = FONT_SANS
    p_status.font.size = Pt(13)
    p_status.font.bold = True
    p_status.font.color.rgb = MUTED_GRAY
    
    # Accent container card for publication screenshots
    add_card(slide18, Inches(7.2), Inches(1.6), Inches(5.3), Inches(4.8), "", [], accent_color=SLATE)
    
    # Insert journal title image (Left side of card)
    journal_img_path = r"figures\journal_title.png"
    if os.path.exists(journal_img_path):
        slide18.shapes.add_picture(journal_img_path, Inches(7.35), Inches(2.3), width=Inches(2.4), height=Inches(3.4))
    else:
        txBox_e = slide18.shapes.add_textbox(Inches(7.35), Inches(2.3), Inches(2.4), Inches(3.4))
        tf_e = txBox_e.text_frame
        tf_e.word_wrap = True
        p_e = tf_e.paragraphs[0]
        p_e.text = "[journal_title.png placeholder]"
        p_e.font.name = FONT_SANS
        p_e.font.size = Pt(12)
        p_e.font.bold = True
        p_e.font.color.rgb = SLATE
        p_e.alignment = PP_ALIGN.CENTER

    # Insert publication screen image (Right side of card)
    pub_img_path = r"figures\publication_screen.png"
    if os.path.exists(pub_img_path):
        slide18.shapes.add_picture(pub_img_path, Inches(9.9), Inches(2.3), width=Inches(2.4), height=Inches(3.4))
    else:
        txBox_e = slide18.shapes.add_textbox(Inches(9.9), Inches(2.3), Inches(2.4), Inches(3.4))
        tf_e = txBox_e.text_frame
        tf_e.word_wrap = True
        p_e = tf_e.paragraphs[0]
        p_e.text = "[publication_screen.png placeholder]"
        p_e.font.name = FONT_SANS
        p_e.font.size = Pt(12)
        p_e.font.bold = True
        p_e.font.color.rgb = SLATE
        p_e.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 19: Conclusion
    # -------------------------------------------------------------
    slide19 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide19)
    add_slide_header(slide19, "Conclusion", 23, TOTAL_SLIDES)
    
    conclusion_points = [
        "1. Implemented zero-phase Butterworth filter and Z-score scaling to remove noise.",
        "2. Built 1D ResNet-18 baseline, achieving a strong 0.7285 Macro F1-score.",
        "3. Developed ST-ReGE GNN with learnable adjacency, capturing lead spatial correlations (F1: 0.7444).",
        "4. Evaluated 1D ViT, showing it underperforms on medium databases due to inductive bias gaps.",
        "5. Built Spatio-Temporal Hybrid combining convolutional morphology with attention rhythm.",
        "6. Calibrated decision boundaries and ensembled models to reach 0.7506 F1 and 0.9344 AUROC.",
        "7. Built async FastAPI prototype and dashboard GUI for real-time triage demo."
    ]
    add_textbox(slide19, Inches(0.8), Inches(1.4), Inches(11.733), Inches(5.4), conclusion_points, size=16, bullet=False)

    # -------------------------------------------------------------
    # SLIDE 19b: Future Directions
    # -------------------------------------------------------------
    slide19b = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide19b)
    add_slide_header(slide19b, "Future Directions", 24, TOTAL_SLIDES)
    
    future_points = [
        "Model Ensembling: Blend 1D ResNet-18, ST-ReGE GNN, and Hybrid models to combine morphological features, spatial relationships, and temporal contexts.",
        "Self-Supervised Pre-training: Investigate self-supervised pre-training (like Masked Auto-Encoders) on massive datasets (e.g., MIMIC-IV-ECG) to improve Transformer performance.",
        "Multi-Modal Integration: Integrate patient demographics (age, sex, clinical history) into the model inputs to build a multi-modal decision support system."
    ]
    add_textbox(slide19b, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.5), future_points, size=16, bullet=True)

    # -------------------------------------------------------------
    # SLIDE 22: Thank You (Redesigned to match Example Slide 17)
    # -------------------------------------------------------------
    slide22 = prs.slides.add_slide(blank_slide_layout)
    apply_bg(slide22)
    
    # Header label
    txBox_cat = slide22.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(8.0), Inches(0.4))
    tf_cat = txBox_cat.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = "THANK YOU"
    p_cat.font.name = FONT_SANS
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = RUST
    
    # Large centered thank you
    txBox_thanks = slide22.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.5))
    tf_thanks = txBox_thanks.text_frame
    tf_thanks.word_wrap = True
    p_th = tf_thanks.paragraphs[0]
    p_th.text = "Thank you\n"
    p_th.font.name = FONT_SERIF
    p_th.font.size = Pt(48)
    p_th.font.bold = True
    p_th.font.color.rgb = SLATE
    p_th.alignment = PP_ALIGN.LEFT
    
    run_th = p_th.add_run()
    run_th.text = "for your attention."
    run_th.font.name = FONT_SERIF
    run_th.font.size = Pt(48)
    run_th.font.bold = True
    run_th.font.italic = True
    run_th.font.color.rgb = RUST
    
    p_welcome = tf_thanks.add_paragraph()
    p_welcome.text = "I welcome your questions and discussion."
    p_welcome.font.name = FONT_SANS
    p_welcome.font.size = Pt(18)
    p_welcome.font.italic = True
    p_welcome.font.color.rgb = SLATE
    p_welcome.space_before = Pt(12)
    p_welcome.alignment = PP_ALIGN.LEFT
    
    # Divider line
    line = slide22.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = SLATE
    line.line.fill.background()
    
    # Candidate details (Left column)
    txBox_cand = slide22.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(5.0), Inches(1.5))
    tf_cand = txBox_cand.text_frame
    tf_cand.word_wrap = True
    tf_cand.margin_left = tf_cand.margin_right = tf_cand.margin_top = tf_cand.margin_bottom = 0
    p_c_label = tf_cand.paragraphs[0]
    p_c_label.text = "CANDIDATE"
    p_c_label.font.name = FONT_SANS
    p_c_label.font.size = Pt(10)
    p_c_label.font.bold = True
    p_c_label.font.color.rgb = MUTED_GRAY
    p_c_label.space_after = Pt(2)
    
    p_c_val = tf_cand.add_paragraph()
    p_c_val.text = "Askar Anafin"
    p_c_val.font.name = FONT_SERIF
    p_c_val.font.size = Pt(16)
    p_c_val.font.bold = True
    p_c_val.font.color.rgb = SLATE
    
    # Supervisor details (Right column)
    txBox_sup = slide22.shapes.add_textbox(Inches(6.8), Inches(5.2), Inches(5.0), Inches(1.5))
    tf_sup = txBox_sup.text_frame
    tf_sup.word_wrap = True
    tf_sup.margin_left = tf_sup.margin_right = tf_sup.margin_top = tf_sup.margin_bottom = 0
    p_s_label = tf_sup.paragraphs[0]
    p_s_label.text = "SCIENTIFIC SUPERVISOR"
    p_s_label.font.name = FONT_SANS
    p_s_label.font.size = Pt(10)
    p_s_label.font.bold = True
    p_s_label.font.color.rgb = MUTED_GRAY
    p_s_label.space_after = Pt(2)
    p_s_val = tf_sup.add_paragraph()
    p_s_val.text = "Minsoo Hahn, PhD"
    p_s_val.font.name = FONT_SERIF
    p_s_val.font.size = Pt(16)
    p_s_val.font.bold = True
    p_s_val.font.color.rgb = SLATE

    # Save presentation
    output_path = r"Presentation\thesis_defense\presentation.pptx"
    try:
        prs.save(output_path)
        print(f"Presentation saved successfully to {output_path}")
    except PermissionError:
        alternative_path = r"Presentation\thesis_defense\presentation_new.pptx"
        try:
            prs.save(alternative_path)
            print(f"Permission denied on {output_path} (likely open in PowerPoint). Saved instead to {alternative_path}")
        except PermissionError:
            # Try versioned fallbacks
            saved = False
            for i in range(1, 100):
                versioned_path = f"Presentation\\thesis_defense\\presentation_v{i}.pptx"
                try:
                    prs.save(versioned_path)
                    print(f"Permission denied on {output_path} and {alternative_path}. Saved instead to {versioned_path}")
                    saved = True
                    break
                except PermissionError:
                    continue
            if not saved:
                print("Error: Could not save presentation to any alternative path due to permission locks.")

if __name__ == "__main__":
    create_presentation()
